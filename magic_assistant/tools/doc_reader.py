from pydantic import BaseModel
import io
import pptx
import docx
from loguru import logger
import fitz


class PdfParser():
    def process(self, file_name: str, file_bytes: bytes) -> str:
        text_data = ""
        doc = fitz.open(stream=file_bytes)
        page_count = doc.page_count

        for page_num in range(page_count):
            text = doc.get_page_text(page_num)

            text = text.strip(" ")
            if len(text) > 0:
                text_data += text + "\f"

        text_data = self.file_clean_text(text_data)
        logger.debug("process %s suc, txt_len:%d" % (file_name, len(text_data)))
        return text_data


class PptxParser():
    def process(self, file_name: str, file_bytes: bytes) -> str:
        text_data = ""
        try:
            data = io.BytesIO(file_bytes)
            ppt = pptx.Presentation(data)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_data += shape.text

            logger.debug("process %s suc, txt_len:%d" % (file_name, len(text_data)))
            return text_data
        except Exception as error:
            logger.error("catch exception: %s", str(error))
            return ""
        finally:
            data.close()


class DocxParser():
    def process(self, file_name: str, file_bytes: bytes) -> str:
        text_data = ""
        try:
            data = io.BytesIO(file_bytes)
            doc = docx.Document(data)
            for para in doc.paragraphs:
                text_data += para.text

            logger.debug("process %s suc, txt_len:%d" % (file_name, len(text_data)))
            return text_data
        except Exception as error:
            logger.error("catch exception: %s" % str(error))
            return ""
        finally:
            data.close()


class DocReader(BaseModel):
    _pdf_parser: PdfParser = PdfParser()
    _docx_parser: DocxParser = DocxParser()
    _pptx_parser: PptxParser = PptxParser()

    def process(self, file_name: str, file_bytes: bytes) -> str:
        if file_name.endswith(".pdf"):
            text_data = self._pdf_parser.process(file_name, file_bytes)
        elif file_name.endswith(".docx"):
            text_data = self._docx_parser.process(file_name, file_bytes)
        elif file_name.endswith(".pptx"):
            text_data = self._pptx_parser.process(file_name, file_bytes)
        else:
            return ""

        return text_data
